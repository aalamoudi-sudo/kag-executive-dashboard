#!/usr/bin/env python3
"""
generate_report.py v3 — مولّد تقارير PPTX
يستخدم القوالب الأصلية ويعبّئ البيانات في أماكنها الصحيحة بالضبط
"""
import sys, json, io, os
from datetime import datetime
from pptx import Presentation

TEMPLATES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")
REPORT_ENGINE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "report_engine")

ORIGINALS = {
    "comprehensive": "comprehensive.pptx",
    "أ": "track-a.pptx",
    "ب": "track-b.pptx",
    "ج": "track-c.pptx",
    "د": "track-d.pptx",
    # التقارير الأربعة الجديدة — تستخدم المولّد الديناميكي
    "daily_ops":  None,
    "executive":  None,
    "approvals":  None,
    "evidence":   None,
}

# ============================================================
# مساعدات
# ============================================================
def today_str():
    return datetime.now().strftime("%Y-%m-%d")

def week_num():
    return datetime.now().isocalendar()[1]

def fmt_date(d):
    if not d: return "—"
    try: return datetime.strptime(str(d)[:10], "%Y-%m-%d").strftime("%Y/%m/%d")
    except: return str(d)[:10]

def status_ar(s):
    if not s: return "—"
    if s in ["مكتملة","مكتمل","معتمدة","معتمد","ضمن المسار"]: return "أخضر ✓"
    if s in ["قيد التنفيذ","تحت المتابعة"]: return "أصفر"
    if s in ["معرضة للخطر","معرض للخطر","متأخرة"]: return "أحمر ✗"
    return s

def is_done(i): return i.get("status","") in ["مكتملة","مكتمل","معتمدة","معتمد"]
def is_active(i): return i.get("status","") == "قيد التنفيذ"
def is_risk(i): return i.get("type","") in ["مخاطرة","مخاطر","risks"]

def get(lst, idx, key="title", default="—"):
    return lst[idx].get(key, default) if idx < len(lst) else default

def bullet_lines(lst, key="title"):
    if not lst: return "لا يوجد"
    return "\n".join(f"• {i.get(key,'')}" for i in lst)

def F(slide, name, text):
    """ابحث عن shape باسمه وعيّن نصه مع الحفاظ على التنسيق"""
    for shape in slide.shapes:
        if shape.name == name and shape.has_text_frame:
            tf = shape.text_frame
            if not tf.paragraphs: return
            # احذف paragraphs الزائدة
            while len(tf.paragraphs) > 1:
                tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
            para = tf.paragraphs[0]
            # احذف runs الزائدة
            while len(para.runs) > 1:
                para.runs[-1]._r.getparent().remove(para.runs[-1]._r)
            # عيّن النص
            t = str(text) if text else "—"
            if para.runs:
                para.runs[0].text = t
            else:
                para.text = t
            return

# ============================================================
# التقرير الشامل
# ============================================================
def fill_comprehensive(prs, state):
    tracks = state.get("tracks", [])
    items  = state.get("items", [])
    risks  = [i for i in items if is_risk(i) and i.get("status") != "مغلقة"]
    done   = [i for i in items if is_done(i)]
    active = [i for i in items if is_active(i)]

    overall = round(sum(t.get("progress",0) for t in tracks)/len(tracks)) if tracks else 0
    worst = "أحمر" if any(t.get("status","") in ["معرض للخطر","معرضة للخطر"] for t in tracks) \
        else "أصفر" if any(t.get("status","") == "تحت المتابعة" for t in tracks) else "أخضر"

    def gt(tid): return next((t for t in tracks if t.get("track")==tid), {})

    # شريحة 1 — الغلاف
    s1 = prs.slides[0]
    F(s1, "Text 4", f"الأسبوع {week_num()} ٢٠٢٦  |  {today_str()}")

    # شريحة 2 — الملخص التنفيذي
    s2 = prs.slides[1]
    F(s2, "Text 17", f"الحالة: {worst}  |  نسبة إنجاز اليوم: {overall}٪")
    F(s2, "Text 21", bullet_lines(done[:3]))
    F(s2, "Text 25", bullet_lines(active[:3]))
    # القضايا الحرجة — 3 صفوف
    risk_rows = [
        ("Text 44","Text 46","Text 48","Text 50","Shape 51"),
        ("Text 54","Text 56","Text 58","Text 60","Text 62"),
        ("Text 64","Text 66","Text 68","Text 70","Text 72"),
    ]
    for idx, (st, sm, sa, sar, sq) in enumerate(risk_rows):
        r = risks[idx] if idx < len(risks) else None
        F(s2, st,  fmt_date(r.get("due",""))  if r else "—")
        F(s2, sm,  r.get("owner","—")         if r else "—")
        F(s2, sa,  "متابعة عاجلة"              if r else "—")
        F(s2, sar, "تأثير على الجدول"          if r else "—")
        F(s2, sq,  r.get("title","—")          if r else "—")

    # شريحة 3 — المسارات الأربعة
    s3 = prs.slides[2]
    track_map = [
        ("أ","Text 32","Text 30","Text 28","Text 26","Text 24"),
        ("ب","Text 44","Text 42","Text 40","Text 38","Text 36"),
        ("ج","Text 56","Text 54","Text 52","Text 50","Text 48"),
        ("د","Text 68","Text 66","Text 64","Text 62","Text 60"),
    ]
    for tid, amsl, alyom, ghad, hal, da3m in track_map:
        t = gt(tid)
        ti = [i for i in items if i.get("track")==tid]
        td = [i for i in ti if is_done(i)]
        ta = [i for i in ti if is_active(i)]
        tn = sorted([i for i in ti if not is_risk(i) and i.get("due","")>today_str()], key=lambda x:x.get("due",""))
        tr = [i for i in ti if is_risk(i) and i.get("status")!="مغلقة"]
        F(s3, amsl,  get(td,0) if td else "—")
        F(s3, alyom, get(ta,0) if ta else "—")
        F(s3, ghad,  get(tn,0) if tn else "—")
        F(s3, hal,   status_ar(t.get("status","")))
        F(s3, da3m,  f"{len(tr)} مخاطر" if tr else "لا يوجد")

    # شريحة 4 — السلامة (نتركها كما هي)

    # شريحة 5 — المخاطر والقرارات
    s5 = prs.slides[4]
    red = [r for r in risks if r.get("status","") in ["معرضة للخطر","معرض للخطر"]]
    yel = [r for r in risks if r.get("status","") in ["تحت المتابعة","قيد التنفيذ"]]
    grn = [r for r in risks if r not in red and r not in yel]
    buckets = [
        (red, "Text 28","Text 22","Text 24","Text 26","Text 30"),
        (yel, "Text 38","Text 32","Text 34","Text 36","Text 40"),
        (grn, "Text 48","Text 42","Text 44","Text 46","Text 50"),
    ]
    for lst, sq, sm, st, stow, _ in buckets:
        r = lst[0] if lst else None
        F(s5, sq,   r.get("title","—")       if r else "—")
        F(s5, sm,   r.get("owner","—")        if r else "—")
        F(s5, st,   fmt_date(r.get("due","")) if r else "—")
        F(s5, stow, "متابعة عاجلة"            if r else "—")
    dec_rows = [
        ("Text 69","Text 67","Text 65","Text 63"),
        ("Text 77","Text 75","Text 73","Text 71"),
        ("Text 85","Text 83","Text 81","Text 79"),
    ]
    for idx, (sw, sa, sar, sr) in enumerate(dec_rows):
        r = risks[idx] if idx < len(risks) else None
        F(s5, sw,  r.get("title","—") if r else "—")
        F(s5, sa,  "تأثير على الجدول" if r else "—")
        F(s5, sar, "متابعة عاجلة"     if r else "—")
        F(s5, sr,  f"خطر-{idx+1}"     if r else "—")

    # شريحة 6 — الجدول الزمني
    s6 = prs.slides[5]
    all_tasks = [i for i in items if not is_risk(i)]
    td = [i for i in all_tasks if is_done(i)][:4]
    ta = [i for i in all_tasks if is_active(i)][:4]
    tn = sorted([i for i in all_tasks if i.get("due","")>today_str()], key=lambda x:x.get("due",""))[:4]
    F(s6, "Text 42", bullet_lines(td))
    F(s6, "Text 31", bullet_lines(ta))
    F(s6, "Text 20", bullet_lines(tn))

# ============================================================
# تقارير المسارات (أ، ب، ج، د — نفس البنية)
# ============================================================
def fill_track(prs, track_key, state):
    tracks = state.get("tracks", [])
    items  = state.get("items", [])
    track  = next((t for t in tracks if t.get("track")==track_key), {})
    ti     = [i for i in items if i.get("track")==track_key]
    risks  = [i for i in ti if is_risk(i) and i.get("status")!="مغلقة"]
    done   = [i for i in ti if is_done(i)]
    active = [i for i in ti if is_active(i)]
    tasks  = [i for i in ti if not is_risk(i)]
    upcoming = sorted([i for i in ti if not is_risk(i) and i.get("due","")>today_str()], key=lambda x:x.get("due",""))
    progress = track.get("progress", 0)

    # شريحة 1 — الغلاف: لا تعديل (اسم المسار موجود في القالب)

    # شريحة 2 — الملخص اليومي
    s2 = prs.slides[1]
    F(s2, "Text 12", f"الحالة المختارة: {status_ar(track.get('status',''))}  |  نسبة إنجاز اليوم: {progress}٪")
    # ملخص الإنجاز — 3 عناصر
    F(s2, "Text 15",
      f"{get(done,0)}\n{get(done,1)}\n{get(done,2)}")
    # الدعم المطلوب
    F(s2, "Text 18",
      f"{get(risks,0)}  |  آخر موعد: {fmt_date(risks[0].get('due','')) if risks else '—'}"
      if risks else "لا يوجد دعم عاجل مطلوب")
    # التحديث التفصيلي — 4 صفوف × 3 أعمدة
    detail_rows = [
        ("Text 42","Text 40","Text 38", 0),  # أنشطة رئيسية
        ("Text 54","Text 52","Text 50", 1),  # مخرجات/اعتمادات
        ("Text 66","Text 64","Text 62", 2),  # تنسيق مع جهات
        ("Text 78","Text 76","Text 74", 3),  # موردين/أطراف
    ]
    for amsl, alyom, ghad, idx in detail_rows:
        F(s2, amsl,  get(done,    idx) if idx < len(done)     else "—")
        F(s2, alyom, get(active,  idx) if idx < len(active)   else "—")
        F(s2, ghad,  get(upcoming,idx) if idx < len(upcoming) else "—")

    # شريحة 3 — تفاصيل الأنشطة (8 صفوف)
    s3 = prs.slides[2]
    task_rows = [
        ("Text 35","Text 33","Text 31","Text 29","Text 27","Text 25"),
        ("Text 51","Text 49","Text 47","Text 45","Text 43","Text 41"),
        ("Text 67","Text 65","Text 63","Text 61","Text 59","Text 57"),
        ("Text 83","Text 81","Text 79","Text 77","Text 75","Text 73"),
        ("Text 99","Text 97","Text 95","Text 93","Text 91","Text 89"),
        ("Text 115","Text 113","Text 111","Text 109","Text 107","Text 105"),
        ("Text 131","Text 129","Text 127","Text 125","Text 123","Text 121"),
        ("Text 147","Text 145","Text 143","Text 141","Text 139","Text 137"),
    ]
    for idx, (stitle, sowner, sstart, send, sstatus, spct) in enumerate(task_rows):
        t = tasks[idx] if idx < len(tasks) else None
        F(s3, stitle,  t.get("title","—")             if t else "—")
        F(s3, sowner,  t.get("owner","—")              if t else "—")
        F(s3, sstart,  fmt_date(t.get("due",""))       if t else "—")
        F(s3, send,    fmt_date(t.get("due",""))        if t else "—")
        F(s3, sstatus, status_ar(t.get("status",""))   if t else "—")
        F(s3, spct,    f"{t.get('progress',0)}٪"        if t else "—")

    # شريحة 4 — المخاطر والقرارات
    s4 = prs.slides[3]
    red = [r for r in risks if r.get("status","") in ["معرضة للخطر","معرض للخطر"]]
    yel = [r for r in risks if r.get("status","") in ["تحت المتابعة","قيد التنفيذ"]]
    grn = [r for r in risks if r not in red and r not in yel]
    risk_buckets = [
        (red, "Text 25","Text 19","Text 21","Text 23","Text 17"),
        (yel, "Text 37","Text 29","Text 31","Text 33","Text 35"),  # fixed: أص row
        (grn, "Text 49","Text 41","Text 43","Text 45","Text 47"),  # fixed: أ row
    ]
    for lst, sq, stime1, sm, stow, stime2 in risk_buckets:
        r = lst[0] if lst else None
        F(s4, sq,     r.get("title","—")       if r else "—")
        F(s4, sm,     r.get("owner","—")        if r else "—")
        F(s4, stow,   "متابعة عاجلة"            if r else "—")
        F(s4, stime1, fmt_date(r.get("due","")) if r else "—")
        F(s4, stime2, fmt_date(r.get("due","")) if r else "—")
    dec_rows4 = [
        ("Text 71","Text 69","Text 67","Text 65"),
        ("Text 79","Text 77","Text 75","Text 73"),
        ("Text 87","Text 85","Text 83","Text 81"),
    ]
    for idx, (sw, sa, sar, sr) in enumerate(dec_rows4):
        r = risks[idx] if idx < len(risks) else None
        F(s4, sw,  r.get("title","—") if r else "—")
        F(s4, sa,  "تأثير على الجدول" if r else "—")
        F(s4, sar, "متابعة"            if r else "—")
        F(s4, sr,  f"خطر-{idx+1}"     if r else "—")

    # شريحة 5 — السلامة: نتركها كما هي

    # شريحة 6 — الجدول الزمني
    s6 = prs.slides[5]
    td6 = [i for i in tasks if is_done(i)][:4]
    ta6 = [i for i in tasks if is_active(i)][:4]
    tn6 = upcoming[:4]
    F(s6, "Text 33", bullet_lines(td6))
    F(s6, "Text 24", bullet_lines(ta6))
    F(s6, "Text 15", bullet_lines(tn6))

    # شريحة 7 — التوصيات: نتركها فارغة للتعبئة اليدوية


# ============================================================
# مولّد احتياطي بدون قوالب — يضمن أن زر التقرير لا يفشل إذا لم تُرفق templates
# ============================================================
def generate_fallback_report(report_type, state):
    from pptx.util import Inches, Pt
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    tracks = state.get("tracks", []) or []
    items = state.get("items", []) or []
    risks = [i for i in items if i.get("type") == "risks" and i.get("status") != "مغلقة"]
    tasks = [i for i in items if i.get("type") == "tasks"]
    done = [i for i in tasks if is_done(i)]
    active = [i for i in tasks if i.get("status") in ["قيد التنفيذ", "تحت المتابعة"]]
    overall = round(sum(t.get("progress",0) for t in tracks)/len(tracks)) if tracks else 0

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid(); bg.fore_color.rgb = __import__('pptx').dml.color.RGBColor(13, 27, 42)
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.8), Inches(1.2))
        tf = box.text_frame; tf.text = title
        tf.paragraphs[0].font.size = Pt(34); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = __import__('pptx').dml.color.RGBColor(217,184,108)
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.15), Inches(11.8), Inches(0.8))
        sub.text_frame.text = subtitle
        sub.text_frame.paragraphs[0].font.size = Pt(18)
        sub.text_frame.paragraphs[0].font.color.rgb = __import__('pptx').dml.color.RGBColor(234,240,247)
        return slide

    def add_bullets(title, lines):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = __import__('pptx').dml.color.RGBColor(13,27,42)
        h = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.55))
        h.text_frame.text = title
        h.text_frame.paragraphs[0].font.size = Pt(26); h.text_frame.paragraphs[0].font.bold = True
        h.text_frame.paragraphs[0].font.color.rgb = __import__('pptx').dml.color.RGBColor(217,184,108)
        body = slide.shapes.add_textbox(Inches(0.9), Inches(1.25), Inches(11.7), Inches(5.7))
        tf = body.text_frame; tf.word_wrap = True; tf.text = ""
        for idx, line in enumerate(lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = "• " + str(line)
            p.font.size = Pt(16)
            p.font.color.rgb = __import__('pptx').dml.color.RGBColor(234,240,247)
        return slide

    add_title_slide("تقرير منصة التحليل التشغيلي", f"نوع التقرير: {report_type} | التاريخ: {today_str()} | التقدم العام: {overall}٪")
    add_bullets("ملخص الأداء", [
        f"عدد المسارات: {len(tracks)}",
        f"إجمالي المهام: {len(tasks)}",
        f"المهام المنجزة: {len(done)}",
        f"المهام تحت التنفيذ/المتابعة: {len(active)}",
        f"المخاطر المفتوحة: {len(risks)}",
    ])
    add_bullets("حالة المسارات", [f"{t.get('id','')} - {t.get('name','')}: {t.get('progress',0)}٪ / {t.get('status','-')}" for t in tracks] or ["لا توجد مسارات"])
    add_bullets("أهم المخاطر والإجراءات", [f"{r.get('title','-')} | المالك: {r.get('owner','-')} | الموعد: {fmt_date(r.get('due',''))}" for r in risks[:8]] or ["لا توجد مخاطر مفتوحة"])
    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()

# ============================================================
# الدالة الرئيسية
# ============================================================
def generate_report(report_type, state):
    # التقارير الأربعة الجديدة — يتم توليدها ديناميكياً بالبيانات الحية
    if report_type in ("daily_ops", "executive", "approvals", "evidence"):
        return generate_dynamic_report(report_type, state)

    tpl_file = ORIGINALS.get(report_type)
    if not tpl_file:
        raise ValueError(f"نوع تقرير غير معروف: {report_type}")

    # ابحث عن القالب في templates/ أولاً، ثم report_engine/ احتياطياً
    tpl_path = os.path.join(TEMPLATES_DIR, tpl_file)
    if not os.path.exists(tpl_path):
        # بدائل في report_engine (daily_comprehensive.pptx)
        fallback_candidates = [
            os.path.join(REPORT_ENGINE_DIR, "daily_comprehensive.pptx"),
            os.path.join(REPORT_ENGINE_DIR, "out_comprehensive.pptx"),
        ]
        for fb in fallback_candidates:
            if os.path.exists(fb):
                tpl_path = fb
                break
        else:
            return generate_fallback_report(report_type, state)

    prs = Presentation(tpl_path)
    if report_type == "comprehensive":
        fill_comprehensive(prs, state)
    else:
        fill_track(prs, report_type, state)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ============================================================
# مولّد ديناميكي للتقارير الأربعة — يسحب البيانات الحية مباشرة
# ============================================================
def generate_dynamic_report(report_type, state):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    tracks = state.get("tracks", []) or []
    items  = list(state.get("items",  []) or [])
    management = state.get("management", {}) or {}
    notifications = state.get("notifications", []) or []
    data_quality = state.get("dataQuality", {}) or {}
    op_summary = state.get("operationalSummary", {}) or {}

    # دمج بيانات مراكز التشغيل داخل التقرير النهائي حتى تظهر كل بيانات النظام
    for a in management.get("actions", []) or []:
        items.append({"type":"tasks", "title":a.get("title"), "track":a.get("track"), "owner":a.get("owner"), "priority":a.get("priority"), "due":a.get("due"), "status":a.get("status"), "evidenceUrl":a.get("evidenceUrl")})
    for ap in management.get("approvals", []) or []:
        items.append({"type":"approval", "title":ap.get("title"), "track":ap.get("track"), "owner":ap.get("owner"), "due":ap.get("due"), "status":ap.get("status"), "impact":ap.get("impact")})
    for ev in management.get("fieldEvidence", []) or []:
        items.append({"type":"evidence", "title":ev.get("title"), "track":ev.get("track"), "owner":ev.get("zone"), "due":ev.get("date"), "status":ev.get("status")})

    risks  = [i for i in items if i.get("type") in ("risks","مخاطرة","مخاطر") and i.get("status") != "مغلقة"]
    tasks  = [i for i in items if i.get("type") == "tasks"]
    done   = [i for i in tasks  if is_done(i)]
    active = [i for i in tasks  if is_active(i)]
    overdue= [i for i in tasks  if i.get("due","") and i.get("due","") < today_str() and not is_done(i)]
    approvals_live = management.get("approvals", []) or []
    actions_live = management.get("actions", []) or []
    evidence_live = management.get("fieldEvidence", []) or []
    meetings_live = management.get("meetings", []) or []
    zones_live = management.get("zones", []) or []
    overall= round(sum(t.get("progress",0) for t in tracks)/len(tracks)) if tracks else 0

    # ألوان الهوية
    C_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # أزرق داكن
    C_GOLD  = RGBColor(0xD9, 0xB8, 0x6C)   # ذهبي
    C_WHITE = RGBColor(0xEA, 0xF0, 0xF7)   # أبيض مائل
    C_GREEN = RGBColor(0x2E, 0xCC, 0x71)
    C_RED   = RGBColor(0xE7, 0x4C, 0x3C)
    C_YEL   = RGBColor(0xF3, 0x9C, 0x12)
    C_GRAY  = RGBColor(0x95, 0xA5, 0xA6)

    def add_slide_bg(color=None):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        bg = sl.background.fill
        bg.solid()
        bg.fore_color.rgb = color or C_BG
        return sl

    def tb(sl, l, t, w, h, text, size=14, bold=False, color=None, align=PP_ALIGN.RIGHT):
        box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf  = box.text_frame; tf.word_wrap = True
        p   = tf.paragraphs[0]; p.text = str(text); p.alignment = align
        p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color or C_WHITE
        return box

    def section_title(sl, text, top=0.42):
        tb(sl, 0.6, top, 12.1, 0.5, text, size=22, bold=True, color=C_GOLD)
        # خط فاصل
        from pptx.util import Emu
        ln = sl.shapes.add_connector(
            __import__('pptx.enum.shapes', fromlist=['MSO_CONNECTOR']).MSO_CONNECTOR.STRAIGHT,
            Inches(0.6), Inches(top+0.52), Inches(12.7), Inches(top+0.52))
        ln.line.color.rgb = C_GOLD; ln.line.width = Emu(12700)

    def badge(sl, l, t, text, color):
        box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(1.6), Inches(0.38))
        tf  = box.text_frame; p = tf.paragraphs[0]
        p.text = text; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11); p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
        box.fill.solid(); box.fill.fore_color.rgb = color
        return box

    def kpi_row(sl, metrics, top=1.1):
        """صف مؤشرات KPI — قائمة من (label, value, color)"""
        n   = len(metrics); w = 12.1/n
        for i,(lbl,val,col) in enumerate(metrics):
            x = 0.6 + i*w
            box = sl.shapes.add_textbox(Inches(x), Inches(top), Inches(w-0.1), Inches(0.9))
            box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x16,0x2A,0x3D)
            tf = box.text_frame
            p0 = tf.paragraphs[0]; p0.text = str(val); p0.alignment = PP_ALIGN.CENTER
            p0.font.size = Pt(26); p0.font.bold = True; p0.font.color.rgb = col
            p1 = tf.add_paragraph(); p1.text = lbl; p1.alignment = PP_ALIGN.CENTER
            p1.font.size = Pt(10); p1.font.color.rgb = C_GRAY

    def data_table(sl, headers, rows, top, height_per_row=0.32):
        """جدول بيانات بسيط"""
        ncol = len(headers); nrow = len(rows)
        col_w = 12.1/ncol
        # رأس الجدول
        for ci,h in enumerate(headers):
            x = 0.6+ci*col_w
            box = sl.shapes.add_textbox(Inches(x), Inches(top), Inches(col_w-0.05), Inches(0.32))
            box.fill.solid(); box.fill.fore_color.rgb = C_GOLD
            tf=box.text_frame; p=tf.paragraphs[0]; p.text=h; p.alignment=PP_ALIGN.CENTER
            p.font.size=Pt(11); p.font.bold=True; p.font.color.rgb=C_BG
        # صفوف البيانات
        for ri,row in enumerate(rows[:12]):
            row_bg = RGBColor(0x12,0x22,0x32) if ri%2 else RGBColor(0x16,0x2A,0x3D)
            for ci,cell in enumerate(row):
                x = 0.6+ci*col_w
                box = sl.shapes.add_textbox(Inches(x), Inches(top+0.35+ri*height_per_row),
                                            Inches(col_w-0.05), Inches(height_per_row))
                box.fill.solid(); box.fill.fore_color.rgb = row_bg
                tf=box.text_frame; p=tf.paragraphs[0]
                p.text=str(cell)[:60]; p.alignment=PP_ALIGN.RIGHT
                p.font.size=Pt(10); p.font.color.rgb=C_WHITE

    # ==========================================
    # تقرير غرفة العمليات اليومي
    # ==========================================
    if report_type == "daily_ops":
        title = "تقرير غرفة العمليات اليومي"

        # شريحة 1 — غلاف
        s = add_slide_bg()
        tb(s, 0.8, 1.0, 11.8, 1.0, "حدائق الملك عبدالله — غرفة العمليات", size=34, bold=True, color=C_GOLD)
        tb(s, 0.8, 2.1, 11.8, 0.5, f"التقرير اليومي  |  {today_str()}  |  الأسبوع {week_num()}", size=18, color=C_WHITE)
        status_color = C_RED if overall<50 else C_YEL if overall<75 else C_GREEN
        tb(s, 0.8, 2.8, 11.8, 0.5, f"نسبة الإنجاز الإجمالية: {overall}٪", size=20, bold=True, color=status_color)

        # شريحة 2 — مؤشرات اليوم
        s = add_slide_bg()
        section_title(s, "ملخص مؤشرات اليوم")
        kpi_row(s, [
            ("إجمالي المهام", len(tasks), C_WHITE),
            ("المنجزة",        len(done),  C_GREEN),
            ("قيد التنفيذ",    len(active),C_YEL),
            ("المتأخرة",       len(overdue),C_RED),
            ("المخاطر المفتوحة",len(risks), C_RED),
        ], top=1.1)
        section_title(s, "حالة المسارات", top=2.3)
        track_rows = []
        for t in tracks:
            st = t.get("status","—")
            col_st = "🟢" if "ضمن" in st else "🟡" if "متابعة" in st else "🔴"
            track_rows.append([t.get("id",""), t.get("name",""), f"{t.get('progress',0)}٪",
                                f"{col_st} {st}", str(t.get("risk",0))])
        data_table(s, ["المسار","الاسم","التقدم","الحالة","المخاطر"], track_rows, top=2.85)

        # شريحة 3 — المهام المتأخرة والحرجة
        s = add_slide_bg()
        section_title(s, "المهام المتأخرة والحرجة")
        od_rows = [[i.get("title","")[:50], i.get("owner","—"), i.get("track","—"),
                    fmt_date(i.get("due","")), i.get("status","—")] for i in overdue[:10]]
        if not od_rows: od_rows=[["لا توجد مهام متأخرة","—","—","—","—"]]
        data_table(s, ["المهمة","المسؤول","المسار","الموعد","الحالة"], od_rows, top=1.05)

        # شريحة 4 — المخاطر المفتوحة
        s = add_slide_bg()
        section_title(s, "المخاطر المفتوحة")
        risk_rows = [[r.get("title","")[:50], r.get("owner","—"), r.get("track","—"),
                      fmt_date(r.get("due","")), status_ar(r.get("status",""))] for r in risks[:10]]
        if not risk_rows: risk_rows=[["لا توجد مخاطر مفتوحة","—","—","—","—"]]
        data_table(s, ["المخاطرة","المسؤول","المسار","الموعد","الحالة"], risk_rows, top=1.05)

        # شريحة 5 — المهام المنجزة اليوم
        s = add_slide_bg()
        section_title(s, "المهام المنجزة")
        done_rows = [[i.get("title","")[:50], i.get("owner","—"), i.get("track","—"),
                      fmt_date(i.get("due",""))] for i in done[:12]]
        if not done_rows: done_rows=[["لا توجد مهام منجزة","—","—","—"]]
        data_table(s, ["المهمة","المسؤول","المسار","تاريخ الإغلاق"], done_rows, top=1.05)

    # ==========================================
    # تقرير اللجنة التنفيذية
    # ==========================================
    elif report_type == "executive":
        title = "تقرير اللجنة التنفيذية"

        # شريحة 1 — غلاف تنفيذي
        s = add_slide_bg()
        tb(s, 0.8, 0.8, 11.8, 1.0, "حدائق الملك عبدالله", size=38, bold=True, color=C_GOLD)
        tb(s, 0.8, 1.9, 11.8, 0.6, "التقرير التنفيذي — اللجنة العليا", size=24, color=C_WHITE)
        tb(s, 0.8, 2.65, 11.8, 0.5, f"التاريخ: {today_str()}  |  الأسبوع {week_num()}", size=16, color=C_GRAY)
        tb(s, 0.8, 3.3,  11.8, 0.5, f"نسبة الإنجاز الإجمالية: {overall}٪", size=22, bold=True,
           color=C_RED if overall<50 else C_YEL if overall<75 else C_GREEN)

        # شريحة 2 — ملخص تنفيذي
        s = add_slide_bg()
        section_title(s, "الملخص التنفيذي")
        worst_status = "معرض للخطر" if any(t.get("status","") in ["معرض للخطر","معرضة للخطر"] for t in tracks)                        else "تحت المتابعة" if any(t.get("status","") == "تحت المتابعة" for t in tracks) else "ضمن المسار"
        tb(s, 0.7, 1.1, 12.0, 0.5,
           f"الوضع العام للمشروع: {worst_status}  |  المسارات النشطة: {len(tracks)}  |  المخاطر الحرجة: {len([r for r in risks if r.get('status','') in ['معرضة للخطر','معرض للخطر']])}",
           size=14, color=C_WHITE)
        kpi_row(s, [
            ("الإنجاز الكلي",   f"{overall}٪",  C_GREEN if overall>=70 else C_YEL if overall>=45 else C_RED),
            ("مهام منجزة",      len(done),       C_GREEN),
            ("مهام نشطة",       len(active),     C_YEL),
            ("مهام متأخرة",     len(overdue),    C_RED),
            ("مخاطر مفتوحة",   len(risks),      C_RED),
            ("جودة البيانات",   f"{data_quality.get('score','—')}٪", C_GREEN if data_quality.get('score',0)>=85 else C_YEL),
        ], top=1.75)
        tb(s, 0.8, 2.65, 11.7, 0.35, f"التوصية التشغيلية: {op_summary.get('recommendation','لا توجد توصية تشغيلية مسجلة')}", size=11, color=C_WHITE)
        section_title(s, "أداء المسارات", top=3.0)
        t_rows = [[t.get("id",""), t.get("name",""), f"{t.get('progress',0)}٪",
                   t.get("status","—"), str(t.get("risk",0))] for t in tracks]
        data_table(s, ["","المسار","التقدم","الحالة","المخاطر"], t_rows, top=3.55, height_per_row=0.30)

        # شريحة 3 — المخاطر الحرجة والقرارات
        s = add_slide_bg()
        section_title(s, "المخاطر الحرجة والقرارات المطلوبة")
        critical_risks = [r for r in risks if r.get("status","") in ["معرضة للخطر","معرض للخطر"]]
        cr_rows = [[r.get("title","")[:55], r.get("owner","—"), r.get("track","—"),
                    fmt_date(r.get("due","")), "قرار عاجل"] for r in critical_risks[:8]]
        if not cr_rows: cr_rows=[["لا توجد مخاطر حرجة","—","—","—","—"]]
        data_table(s, ["المخاطرة","المسؤول","المسار","الموعد","الإجراء"], cr_rows, top=1.05)

        # شريحة 4 — الجدول الزمني للمشروع
        s = add_slide_bg()
        section_title(s, "الجدول الزمني — ما تم وما هو قادم")
        upcoming = sorted([i for i in tasks if not is_done(i) and i.get("due","")>today_str()],
                          key=lambda x:x.get("due",""))
        sec_rows  = [[i.get("title","")[:50], i.get("track","—"), fmt_date(i.get("due",""))]
                     for i in done[-6:]]
        up_rows   = [[i.get("title","")[:50], i.get("track","—"), fmt_date(i.get("due",""))]
                     for i in upcoming[:6]]
        tb(s, 0.7, 1.05, 5.8, 0.3, "✅ المنجز", size=13, bold=True, color=C_GREEN)
        data_table(s, ["المهمة","المسار","التاريخ"], sec_rows if sec_rows else [["لا يوجد","",""]],
                   top=1.35, height_per_row=0.30)
        tb(s, 6.9, 1.05, 5.8, 0.3, "⏳ القادم", size=13, bold=True, color=C_YEL)
        # جدول ثانٍ في النصف الأيسر
        ncol=3; col_w_r=6.1/ncol
        hdr2=["المهمة","المسار","التاريخ"]
        for ci,h in enumerate(hdr2):
            box=s.shapes.add_textbox(Inches(6.9+ci*col_w_r),Inches(1.35),Inches(col_w_r-0.05),Inches(0.32))
            box.fill.solid(); box.fill.fore_color.rgb=C_GOLD
            tf=box.text_frame; p=tf.paragraphs[0]; p.text=h; p.alignment=PP_ALIGN.CENTER
            p.font.size=Pt(11); p.font.bold=True; p.font.color.rgb=C_BG
        for ri,row in enumerate((up_rows if up_rows else [["لا يوجد","",""]])[:6]):
            row_bg=RGBColor(0x12,0x22,0x32) if ri%2 else RGBColor(0x16,0x2A,0x3D)
            for ci,cell in enumerate(row):
                box=s.shapes.add_textbox(Inches(6.9+ci*col_w_r),Inches(1.7+ri*0.30),Inches(col_w_r-0.05),Inches(0.30))
                box.fill.solid(); box.fill.fore_color.rgb=row_bg
                tf=box.text_frame; p=tf.paragraphs[0]; p.text=str(cell)[:40]; p.alignment=PP_ALIGN.RIGHT
                p.font.size=Pt(10); p.font.color.rgb=C_WHITE

    # ==========================================
    # تقرير الاعتمادات والتصعيد
    # ==========================================
    elif report_type == "approvals":
        title = "تقرير الاعتمادات والتصعيد"
        approval_items = approvals_live or [i for i in items if i.get("type") in ("approval","اعتماد","موافقة") or "اعتماد" in i.get("title","")]
        pending_tasks  = [i for i in tasks if not is_done(i)]
        escalated      = [i for i in items if i.get("status","") in ("معرض للخطر","معرضة للخطر","متأخرة","متأخر")]

        # شريحة 1 — غلاف
        s = add_slide_bg()
        tb(s, 0.8, 1.0, 11.8, 1.0, "تقرير الاعتمادات والتصعيد", size=34, bold=True, color=C_GOLD)
        tb(s, 0.8, 2.1, 11.8, 0.5, f"حدائق الملك عبدالله  |  {today_str()}", size=18, color=C_WHITE)

        # شريحة 2 — ملخص الاعتمادات
        s = add_slide_bg()
        section_title(s, "ملخص الاعتمادات")
        kpi_row(s, [
            ("إجمالي الاعتمادات",  len(approval_items), C_WHITE),
            ("المعلقة",            len([i for i in pending_tasks if not is_done(i)]), C_YEL),
            ("المتأخرة",           len(overdue), C_RED),
            ("حالات التصعيد",      len(escalated), C_RED),
        ], top=1.1)
        section_title(s, "المهام بانتظار اعتماد", top=2.3)
        pend_source = approval_items if approval_items else pending_tasks
        pend_rows = [[i.get("title","")[:50], i.get("owner","—"), i.get("track","—"),
                      fmt_date(i.get("due","")), i.get("status","—")] for i in pend_source[:10]]
        if not pend_rows: pend_rows=[["لا توجد مهام معلقة","—","—","—","—"]]
        data_table(s, ["المهمة","المسؤول","المسار","الموعد","الحالة"], pend_rows, top=2.85)

        # شريحة 3 — حالات التصعيد
        s = add_slide_bg()
        section_title(s, "حالات التصعيد الحرجة")
        esc_rows = [[i.get("title","")[:50], i.get("owner","—"), i.get("track","—"),
                     fmt_date(i.get("due","")), status_ar(i.get("status",""))] for i in escalated[:12]]
        if not esc_rows: esc_rows=[["لا توجد حالات تصعيد","—","—","—","—"]]
        data_table(s, ["العنصر","المسؤول","المسار","الموعد","الحالة"], esc_rows, top=1.05)

        # شريحة 4 — المخاطر المرتبطة بالاعتمادات
        s = add_slide_bg()
        section_title(s, "المخاطر المرتبطة بالاعتمادات المتأخرة")
        risk_rows2 = [[r.get("title","")[:50], r.get("owner","—"), r.get("track","—"),
                       fmt_date(r.get("due","")), "متابعة عاجلة"] for r in risks[:10]]
        if not risk_rows2: risk_rows2=[["لا توجد مخاطر","—","—","—","—"]]
        data_table(s, ["المخاطرة","المسؤول","المسار","الموعد","الإجراء"], risk_rows2, top=1.05)

    # ==========================================
    # تقرير الأدلة الميدانية
    # ==========================================
    elif report_type == "evidence":
        title = "تقرير الأدلة الميدانية"
        ev_items = evidence_live or [i for i in items if i.get("type") in ("evidence","دليل","أدلة")
                    or any(k in i.get("title","") for k in ["دليل","صورة","تحقق","إغلاق","فحص"])]

        # شريحة 1 — غلاف
        s = add_slide_bg()
        tb(s, 0.8, 1.0, 11.8, 1.0, "تقرير الأدلة الميدانية", size=34, bold=True, color=C_GOLD)
        tb(s, 0.8, 2.1, 11.8, 0.5, f"حدائق الملك عبدالله  |  {today_str()}", size=18, color=C_WHITE)
        tb(s, 0.8, 2.75, 11.8, 0.4, f"إجمالي الأدلة: {len(ev_items)}  |  المهام المنجزة: {len(done)}", size=16, color=C_GRAY)

        # شريحة 2 — ملخص الأدلة
        s = add_slide_bg()
        section_title(s, "ملخص الأدلة الميدانية")
        kpi_row(s, [
            ("إجمالي الأدلة",    len(ev_items), C_WHITE),
            ("المهام المنجزة",   len(done),     C_GREEN),
            ("قيد التنفيذ",      len(active),   C_YEL),
            ("مخاطر مفتوحة",    len(risks),    C_RED),
        ], top=1.1)
        section_title(s, "قائمة الأدلة الميدانية", top=2.3)
        ev_rows = [[i.get("title","")[:55], i.get("track","—"), i.get("owner","—"),
                    fmt_date(i.get("due","")), status_ar(i.get("status",""))] for i in (ev_items or done[:10])[:12]]
        if not ev_rows: ev_rows=[["لا توجد أدلة مسجلة","—","—","—","—"]]
        data_table(s, ["الدليل/المهمة","المسار","المسؤول","التاريخ","الحالة"], ev_rows, top=2.85)

        # شريحة 3 — المهام المنجزة كأدلة إغلاق
        s = add_slide_bg()
        section_title(s, "المهام المغلقة — أدلة الإغلاق")
        cl_rows = [[i.get("title","")[:55], i.get("track","—"), i.get("owner","—"),
                    fmt_date(i.get("due",""))] for i in done[:15]]
        if not cl_rows: cl_rows=[["لا توجد مهام منجزة","—","—","—"]]
        data_table(s, ["المهمة","المسار","المسؤول","تاريخ الإغلاق"], cl_rows, top=1.05)

        # شريحة 4 — مهام تحتاج توثيق
        s = add_slide_bg()
        section_title(s, "مهام تحتاج توثيق أو دليل ميداني")
        need_doc = [i for i in active if not is_done(i)][:12]
        nd_rows = [[i.get("title","")[:55], i.get("track","—"), i.get("owner","—"),
                    fmt_date(i.get("due",""))] for i in need_doc]
        if not nd_rows: nd_rows=[["جميع المهام النشطة موثقة","—","—","—"]]
        data_table(s, ["المهمة","المسار","المسؤول","الموعد"], nd_rows, top=1.05)

    else:
        return generate_fallback_report(report_type, state)

    # تذييل على كل شريحة
    for sl in prs.slides:
        foot = sl.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.35))
        tf   = foot.text_frame; p = tf.paragraphs[0]
        p.text = f"حدائق الملك عبدالله  |  {title}  |  {today_str()}  |  سري — للاستخدام الداخلي"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(9); p.font.color.rgb = C_GRAY

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()

if __name__ == "__main__":
    data  = json.loads(sys.stdin.read())
    result = generate_report(data.get("type","comprehensive"), data.get("state",{}))
    sys.stdout.buffer.write(result)
